"""
文档服务：文件解析、文本分块、向量化存储，以及文档的列表与删除。
"""

from __future__ import annotations

import logging
import os
import tempfile
from typing import List
from uuid import uuid4

logger = logging.getLogger(__name__)


class DocumentService:
    """文档上传、解析、分块、向量化及管理服务。"""

    # ------------------------------------------------------------------
    # 6.1 文件解析器
    # ------------------------------------------------------------------

    def parse_file(self, tmp_path: str, filename: str) -> str:
        """
        解析文件为纯文本。

        :param tmp_path: 临时文件路径
        :param filename: 原始文件名（用于判断扩展名）
        :raises ValueError: 不支持的文件格式
        :return: 提取的文本内容
        """
        ext = os.path.splitext(filename)[1].lower()

        if ext == ".pdf":
            return self._parse_pdf(tmp_path)
        elif ext == ".docx":
            return self._parse_docx(tmp_path)
        elif ext == ".pptx":
            return self._parse_pptx(tmp_path)
        elif ext in (".txt", ".md"):
            return self._parse_text(tmp_path)
        else:
            raise ValueError(f"不支持的文件格式：{ext}")

    def _parse_pdf(self, tmp_path: str) -> str:
        import pdfplumber

        pages_text: List[str] = []

        with pdfplumber.open(tmp_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if not text.strip():
                    # 扫描版页面，跳过（不做 OCR，避免卡死）
                    logger.warning("第 %d 页无文字内容，跳过", page_num)
                    text = ""
                pages_text.append(text)

        return "\n".join(pages_text)

    def _parse_docx(self, tmp_path: str) -> str:
        from docx import Document

        doc = Document(tmp_path)
        return "\n".join(para.text for para in doc.paragraphs)

    def _parse_pptx(self, tmp_path: str) -> str:
        from pptx import Presentation

        prs = Presentation(tmp_path)
        texts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        texts.append(para.text)
        return "\n".join(texts)

    def _parse_text(self, tmp_path: str) -> str:
        with open(tmp_path, "r", encoding="utf-8") as f:
            return f.read()

    # ------------------------------------------------------------------
    # 6.2 文本分块
    # ------------------------------------------------------------------

    def chunk_text(self, text: str) -> List[str]:
        """
        按滑动窗口将文本分块。

        :param text: 待分块文本
        :return: 文本块列表
        """
        if not text:
            return []

        from config import get_config
        cfg = get_config()
        chunk_size = cfg.CHUNK_SIZE
        chunk_overlap = cfg.CHUNK_OVERLAP

        chunks: List[str] = []
        start = 0
        text_len = len(text)

        while start < text_len:
            end = start + chunk_size
            chunks.append(text[start:end])
            if end >= text_len:
                break
            start += chunk_size - chunk_overlap

        return chunks

    # ------------------------------------------------------------------
    # 6.4 完整上传流程
    # ------------------------------------------------------------------

    def upload_and_process(
        self,
        file_bytes: bytes,
        filename: str,
        subject_id: int,
        user_id: int,
    ) -> dict:
        """
        完整的文件上传与处理流程。

        :param file_bytes: 文件二进制内容
        :param filename: 原始文件名
        :param subject_id: 所属学科 ID
        :param user_id: 上传用户 ID
        :return: {"success": bool, "doc_id": int, "error": str}
        """
        from database import get_session, Document, Chunk
        from services.embedding_service import EmbeddingService

        tmp_path = os.path.join(tempfile.gettempdir(), f"{uuid4()}_{filename}")
        doc_id: int | None = None

        try:
            # 写入临时文件
            with open(tmp_path, "wb") as f:
                f.write(file_bytes)

            # 1. 写 documents 记录（status='pending'）
            with get_session() as session:
                doc = Document(
                    subject_id=subject_id,
                    user_id=user_id,
                    filename=filename,
                    status="pending",
                )
                session.add(doc)
                session.flush()
                doc_id = doc.id

            # 2. 更新 status='processing'
            self._update_doc_status(doc_id, "processing")

            # 3. 解析文本
            text = self.parse_file(tmp_path, filename)

            # 4. 分块
            chunks = self.chunk_text(text)

            # 5. 生成向量
            embedding_service = EmbeddingService()
            vectors = embedding_service.embed_texts(chunks) if chunks else []

            # 6. 存入 PGVector
            if chunks:
                self._store_vectors(chunks, vectors, doc_id, subject_id, filename)

            # 7. 写 chunks 表
            with get_session() as session:
                for idx, chunk_content in enumerate(chunks):
                    chunk = Chunk(
                        document_id=doc_id,
                        subject_id=subject_id,
                        chunk_index=idx,
                        content=chunk_content,
                    )
                    session.add(chunk)

            # 8. 更新 status='completed'
            self._update_doc_status(doc_id, "completed")

            return {"success": True, "doc_id": doc_id, "error": ""}

        except Exception as e:
            logger.error("文档处理失败：%s", e)
            if doc_id is not None:
                self._update_doc_status(doc_id, "failed", str(e))
            return {"success": False, "doc_id": doc_id, "error": str(e)}

        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)

    def _store_vectors(
        self,
        chunks: List[str],
        vectors: List[List[float]],
        doc_id: int,
        subject_id: int,
        filename: str,
    ) -> None:
        """将文本块及向量存入 PGVector collection。"""
        from langchain_postgres import PGVector
        from langchain_openai import OpenAIEmbeddings
        from langchain_core.documents import Document as LCDocument
        from config import get_config

        cfg = get_config()
        embeddings = OpenAIEmbeddings(
            model=cfg.LLM_EMBEDDING_MODEL,
            openai_api_key=cfg.LLM_API_KEY,
            openai_api_base=cfg.LLM_BASE_URL,
        )

        vector_store = PGVector(
            embeddings=embeddings,
            collection_name=f"subject_{subject_id}",
            connection=cfg.DATABASE_URL,
        )

        docs = [
            LCDocument(
                page_content=chunk,
                metadata={
                    "doc_id": doc_id,
                    "subject_id": subject_id,
                    "filename": filename,
                    "chunk_index": idx,
                },
            )
            for idx, chunk in enumerate(chunks)
        ]

        vector_store.add_documents(docs)

    def _update_doc_status(
        self, doc_id: int, status: str, error: str | None = None
    ) -> None:
        """更新 documents 表中的 status 和 error 字段。"""
        from database import get_session, Document

        with get_session() as session:
            doc = session.get(Document, doc_id)
            if doc:
                doc.status = status
                if error is not None:
                    doc.error = error

    # ------------------------------------------------------------------
    # 6.7 列表与删除
    # ------------------------------------------------------------------

    def list_documents(self, subject_id: int, user_id: int) -> List[dict]:
        """
        查询指定学科下当前用户的所有文档。

        :param subject_id: 学科 ID
        :param user_id: 用户 ID
        :return: 文档信息列表
        """
        from database import get_session, Document

        with get_session() as session:
            docs = (
                session.query(Document)
                .filter(
                    Document.subject_id == subject_id,
                    Document.user_id == user_id,
                )
                .order_by(Document.created_at.desc())
                .all()
            )
            return [
                {
                    "id": doc.id,
                    "filename": doc.filename,
                    "status": doc.status,
                    "error": doc.error,
                    "created_at": doc.created_at,
                }
                for doc in docs
            ]

    def delete_document(
        self, doc_id: int, subject_id: int, user_id: int
    ) -> dict:
        """
        删除文档记录及其对应的 PGVector 向量。

        :param doc_id: 文档 ID
        :param subject_id: 学科 ID（用于定位 collection）
        :param user_id: 用户 ID（验证归属）
        :return: {"success": bool, "error": str}
        """
        from database import get_session, Document

        try:
            # 验证归属并删除数据库记录
            with get_session() as session:
                doc = (
                    session.query(Document)
                    .filter(
                        Document.id == doc_id,
                        Document.subject_id == subject_id,
                        Document.user_id == user_id,
                    )
                    .first()
                )
                if doc is None:
                    return {"success": False, "error": "文档不存在或无权限删除"}
                session.delete(doc)

            # 删除 PGVector 中对应向量
            self._delete_vectors(doc_id, subject_id)

            return {"success": True, "error": ""}

        except Exception as e:
            logger.error("删除文档失败：%s", e)
            return {"success": False, "error": str(e)}

    def _delete_vectors(self, doc_id: int, subject_id: int) -> None:
        """从 PGVector collection 中删除指定 doc_id 的所有向量。"""
        from langchain_postgres import PGVector
        from langchain_openai import OpenAIEmbeddings
        from config import get_config

        cfg = get_config()
        embeddings = OpenAIEmbeddings(
            model=cfg.LLM_EMBEDDING_MODEL,
            openai_api_key=cfg.LLM_API_KEY,
            openai_api_base=cfg.LLM_BASE_URL,
        )

        vector_store = PGVector(
            embeddings=embeddings,
            collection_name=f"subject_{subject_id}",
            connection=cfg.DATABASE_URL,
        )

        # 通过 metadata filter 删除该文档的所有向量
        try:
            vector_store.delete(filter={"doc_id": doc_id})
        except Exception as e:
            logger.warning("删除 PGVector 向量失败（doc_id=%d）：%s", doc_id, e)
